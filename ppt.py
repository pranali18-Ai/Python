from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt
from PIL import ImageFont


def add_text_slide(presentation, text, font_file):
    slide_layout = presentation.slide_layouts[1]  # Use layout for a text slide
    slide = presentation.slides.add_slide(slide_layout)
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = text
    p.font.name = font_file
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.LEFT


def create_presentation(slide1_content_file, slide2_content_file, font_file):
    presentation = Presentation()

    # Slide 1
    with open(slide1_content_file, 'r') as f:
        slide1_content = f.read()
    add_text_slide(presentation, slide1_content, font_file)

    # Slide 2
    with open(slide2_content_file, 'r') as f:
        slide2_content = f.read()
    add_text_slide(presentation, slide2_content, font_file)

    presentation.save('output.pptx')


# Provide the file names for slide content and font
slide1_content_file = 'sample_slide1_input.txt'
slide2_content_file = 'sample_slide2_input.txt'
font_file = 'sample_font_file.ttf'

create_presentation(slide1_content_file, slide2_content_file, font_file)
